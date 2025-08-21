"""Slide reconstruction from image to PPTX using OCR and shape detection.

This module reads an image of a PowerPoint slide, extracts textual and
non-textual elements, and rebuilds an editable slide in a PPTX file. The
implementation relies on open‑source libraries such as OpenCV for element
analysis, pytesseract for OCR, and python-pptx for generating the final slide.

The approach is heuristic and may require manual adjustments depending on the
complexity of the original slide. The code is organised into modular functions
so individual steps can be improved or swapped out for more advanced models.
"""

from __future__ import annotations

import argparse
import dataclasses
from dataclasses import dataclass
from typing import List, Tuple, Optional

from PIL import Image, ImageStat
import numpy as np
import cv2  # type: ignore
import pytesseract  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor


@dataclass
class TextElement:
    text: str
    left: int
    top: int
    width: int
    height: int
    font_size: float
    color: Tuple[int, int, int]


@dataclass
class ShapeElement:
    kind: str  # e.g. "rectangle", "ellipse", "line", "image"
    left: int
    top: int
    width: int
    height: int
    color: Tuple[int, int, int]
    path: Optional[str] = None  # used for images


def extract_text_elements(image: Image.Image) -> List[TextElement]:
    """Extract textual elements using Tesseract OCR.

    Returns a list of TextElement with bounding boxes and colour estimates.
    """
    rgb = image.convert("RGB")
    data = pytesseract.image_to_data(rgb, output_type=pytesseract.Output.DICT)

    elements: List[TextElement] = []
    for i, word in enumerate(data["text"]):
        if not word.strip():
            continue
        left, top, width, height = (
            data["left"][i],
            data["top"][i],
            data["width"][i],
            data["height"][i],
        )
        crop = rgb.crop((left, top, left + width, top + height))
        stat = ImageStat.Stat(crop)
        color = tuple(int(c) for c in stat.mean[:3])
        font_size = height  # crude estimate; true font size requires more work
        elements.append(
            TextElement(
                text=word,
                left=left,
                top=top,
                width=width,
                height=height,
                font_size=font_size,
                color=color,
            )
        )
    return elements


def detect_shapes(image: Image.Image) -> List[ShapeElement]:
    """Detect simple geometric shapes and embedded images using OpenCV.

    This function performs edge detection and contour analysis to detect basic
    shapes like rectangles, lines and circles. More complex elements such as
    charts or tables would require specialised models and are not covered.
    """
    rgb = np.array(image.convert("RGB"))
    gray = cv2.cvtColor(rgb, cv2.COLOR_RGB2GRAY)
    blur = cv2.GaussianBlur(gray, (5, 5), 0)
    _, thresh = cv2.threshold(blur, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    elements: List[ShapeElement] = []
    h, w = gray.shape
    for cnt in contours:
        x, y, width, height = cv2.boundingRect(cnt)
        area = width * height
        if area < 100:  # skip tiny artefacts
            continue
        crop = rgb[y : y + height, x : x + width]
        stat = ImageStat.Stat(Image.fromarray(crop))
        color = tuple(int(c) for c in stat.mean[:3])
        peri = cv2.arcLength(cnt, True)
        approx = cv2.approxPolyDP(cnt, 0.04 * peri, True)
        kind = "rectangle"
        if len(approx) > 8:
            kind = "ellipse"
        elif len(approx) == 2:
            kind = "line"
        elements.append(
            ShapeElement(
                kind=kind,
                left=x,
                top=y,
                width=width,
                height=height,
                color=color,
            )
        )
    return elements


def create_pptx(
    image: Image.Image,
    texts: List[TextElement],
    shapes: List[ShapeElement],
    output_path: str,
) -> None:
    """Build a PPTX slide from detected elements."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set slide dimensions to match input image
    prs.slide_width = Inches(image.width / 96)  # assuming 96 DPI
    prs.slide_height = Inches(image.height / 96)

    for s in shapes:
        left = Inches(s.left / 96)
        top = Inches(s.top / 96)
        width = Inches(s.width / 96)
        height = Inches(s.height / 96)
        if s.kind == "rectangle":
            shape = slide.shapes.add_shape(1, left, top, width, height)
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*s.color)
        elif s.kind == "ellipse":
            shape = slide.shapes.add_shape(9, left, top, width, height)
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*s.color)
        elif s.kind == "line":
            line = slide.shapes.add_line(left, top, left + width, top + height)
            line.line.color.rgb = RGBColor(*s.color)
        elif s.kind == "image" and s.path:
            slide.shapes.add_picture(s.path, left, top, width=width, height=height)

    for t in texts:
        left = Inches(t.left / 96)
        top = Inches(t.top / 96)
        width = Inches(t.width / 96)
        height = Inches(t.height / 96)
        box = slide.shapes.add_textbox(left, top, width, height)
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = t.text
        font = run.font
        font.size = Pt(t.font_size)
        font.color.rgb = RGBColor(*t.color)

    prs.save(output_path)


@dataclasses.dataclass
class SlideElements:
    texts: List[TextElement]
    shapes: List[ShapeElement]


def process_slide(image_path: str) -> SlideElements:
    image = Image.open(image_path)
    texts = extract_text_elements(image)
    shapes = detect_shapes(image)
    return SlideElements(texts=texts, shapes=shapes)


def main() -> None:
    parser = argparse.ArgumentParser(description="Rebuild a PPTX slide from an image")
    parser.add_argument("image", help="Path to slide image (PNG/JPG)")
    parser.add_argument("output", help="Path to output PPTX file")
    args = parser.parse_args()

    elements = process_slide(args.image)
    image = Image.open(args.image)
    create_pptx(image, elements.texts, elements.shapes, args.output)
    print(f"Generated {args.output}")


if __name__ == "__main__":
    main()
